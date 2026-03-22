"""
core/file_handler.py
Geração de arquivos (PDF, DOCX, PPTX, TXT) + leitura de uploads.

Formatos suportados no GENERATE_FILE:
  pdf   → PDF com design visual, logo roxa, dicas e palavras-chave
  docx  → Word document formatado
  pptx  → Apresentação PowerPoint
  txt   → Texto simples
  doc   → alias para docx
"""

import base64
import io
import json
import os
import re
from pathlib import Path

import streamlit as st

from core.database import append_message, cached_load_conversation

PROF_NAME  = os.getenv("PROFESSOR_NAME", "Teacher Tati")
LOGO_PATH  = os.getenv("PROFESSOR_PHOTO", "assets/professor.jpg")


# ══════════════════════════════════════════════════════════════════════════════
# INTERCEPTAÇÃO
# ══════════════════════════════════════════════════════════════════════════════

def intercept_file_generation(reply_text: str, username: str, conv_id: str) -> str:
    try:
        match = re.search(
            r'<<<GENERATE_FILE>>>\s*(\{.*?\})\s*<<<END_FILE>>>',
            reply_text, re.DOTALL,
        )
        if not match:
            append_message(username, conv_id, "assistant", reply_text)
            return reply_text

        meta     = json.loads(match.group(1))
        fmt      = meta.get("format", "pdf").lower().strip(".")
        title    = meta.get("title", "Activity")
        content  = meta.get("content", "")
        filename = meta.get("filename", f"activity.{fmt}")

        # normaliza extensão
        fmt_map  = {"doc": "docx", "ppt": "pptx", "text": "txt"}
        fmt      = fmt_map.get(fmt, fmt)
        if not filename.endswith(f".{fmt}"):
            filename = Path(filename).stem + f".{fmt}"

        out_dir  = Path("data/generated")
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = out_dir / filename

        # gera o arquivo
        generators = {
            "pdf":  _gen_pdf,
            "docx": _gen_docx,
            "pptx": _gen_pptx,
            "txt":  _gen_txt,
        }
        gen = generators.get(fmt, _gen_pdf)
        gen(title, content, out_path)

        mime_map = {
            "pdf":  "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "txt":  "text/plain",
        }

        with open(out_path, "rb") as f:
            file_bytes = f.read()

        st.session_state["_pending_download"] = {
            "b64":      base64.b64encode(file_bytes).decode(),
            "filename": filename,
            "mime":     mime_map.get(fmt, "application/octet-stream"),
        }

        display_msg = (
            f"File generated: **{filename}**\n\n_{title}_\n\n"
            "Click **Download file** below to save it."
        )
        append_message(username, conv_id, "assistant", display_msg, is_file=True)
        cached_load_conversation.clear()
        return display_msg

    except Exception as e:
        err = f"Sorry, I couldn't generate the file: {e}"
        append_message(username, conv_id, "assistant", err)
        cached_load_conversation.clear()
        return err


# ══════════════════════════════════════════════════════════════════════════════
# GERADOR PDF — design profissional
# ══════════════════════════════════════════════════════════════════════════════

def _gen_pdf(title: str, content: str, out_path: Path) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units     import cm
    from reportlab.lib           import colors
    from reportlab.lib.styles    import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums     import TA_CENTER
    from reportlab.platypus      import (
        SimpleDocTemplate, Paragraph, Spacer, HRFlowable,
        Table, TableStyle, Image as RLImage,
    )
    from reportlab.platypus import Flowable
    from PIL import Image as PILImage, ImageDraw as PILDraw

    # ── Paleta ────────────────────────────────────────────────────────────────
    C_PURPLE  = colors.HexColor("#7c3aed")
    C_PURPLE2 = colors.HexColor("#a78bfa")
    C_PURPLE3 = colors.HexColor("#f5f3ff")
    C_GOLD    = colors.HexColor("#f0a500")
    C_GREEN   = colors.HexColor("#059669")
    C_GREEN_L = colors.HexColor("#d1fae5")
    C_BLUE    = colors.HexColor("#1d4ed8")
    C_BLUE_L  = colors.HexColor("#dbeafe")
    C_ORANGE_L= colors.HexColor("#ffedd5")
    C_GRAY    = colors.HexColor("#6b7280")
    C_DARK    = colors.HexColor("#0d1117")
    C_WHITE   = colors.white
    C_BORDER  = colors.HexColor("#e5e7eb")
    C_ANS_BG  = colors.HexColor("#f0fdf4")
    C_TIP_BG  = colors.HexColor("#fffbeb")

    W = A4[0] - 4.0 * cm

    class ColorBar(Flowable):
        def __init__(self, c, h=4, w=None):
            super().__init__(); self._c=c; self._h=h; self._w=w
        def wrap(self, aw, ah): return (self._w or aw), self._h
        def draw(self):
            self.canv.setFillColor(self._c)
            self.canv.rect(0, 0, self._w or self.canv._pagesize[0], self._h, fill=1, stroke=0)

    def _circle_crop(raw: bytes, sz=72) -> bytes:
        img  = PILImage.open(io.BytesIO(raw)).convert("RGBA").resize((sz*2,sz*2), PILImage.LANCZOS)
        mask = PILImage.new("L", img.size, 0)
        PILDraw.Draw(mask).ellipse((0,0)+img.size, fill=255)
        img.putalpha(mask)
        buf = io.BytesIO(); img.save(buf,"PNG"); buf.seek(0); return buf.read()

    def _box(para, bg, border, pv=8, ph=12):
        t = Table([[para]], colWidths=[W])
        t.setStyle(TableStyle([
            ("BACKGROUND",   (0,0),(-1,-1), bg),
            ("BOX",          (0,0),(-1,-1), 0.8, border),
            ("TOPPADDING",   (0,0),(-1,-1), pv),
            ("BOTTOMPADDING",(0,0),(-1,-1), pv),
            ("LEFTPADDING",  (0,0),(-1,-1), ph),
            ("RIGHTPADDING", (0,0),(-1,-1), ph),
            ("ROWBACKGROUNDS",(0,0),(-1,-1),[bg]),
        ]))
        return t

    SS = getSampleStyleSheet()
    def S(n,**kw): return ParagraphStyle(n, parent=SS["Normal"], **kw)

    s_htitle = S("HT", fontSize=20, fontName="Helvetica-Bold", textColor=C_WHITE,   alignment=TA_CENTER, leading=26)
    s_hsub   = S("HS", fontSize=9,  fontName="Helvetica",      textColor=C_PURPLE2, alignment=TA_CENTER)
    s_body   = S("BO", fontSize=11, fontName="Helvetica",      textColor=C_DARK, leading=18, spaceAfter=5)
    s_sec    = S("SE", fontSize=13, fontName="Helvetica-Bold", textColor=C_PURPLE, spaceAfter=4, spaceBefore=8)
    s_ans    = S("AK", fontSize=13, fontName="Helvetica-Bold", textColor=C_GREEN,  spaceAfter=4, spaceBefore=8)
    s_ansi   = S("AI", fontSize=10, fontName="Helvetica",      textColor=colors.HexColor("#065f46"), leading=16)
    s_tip    = S("TI", fontSize=10, fontName="Helvetica-Oblique", textColor=colors.HexColor("#92400e"), leading=16)
    s_gram   = S("GR", fontSize=10.5, fontName="Helvetica",    textColor=colors.HexColor("#1e3a5f"), leading=17)
    s_exnum  = S("EN", fontSize=11, fontName="Helvetica-Bold", textColor=C_DARK, leading=18, leftIndent=4)
    s_foot   = S("FO", fontSize=8,  fontName="Helvetica",      textColor=C_GRAY, alignment=TA_CENTER)
    s_kw     = S("KW", fontSize=9.5, fontName="Helvetica",     textColor=colors.HexColor("#4c1d95"), leading=16)
    s_exh    = S("EH", fontSize=12, fontName="Helvetica-Bold", textColor=C_WHITE)

    doc   = SimpleDocTemplate(str(out_path), pagesize=A4,
                              leftMargin=2*cm, rightMargin=2*cm,
                              topMargin=2*cm,  bottomMargin=2*cm)
    story = []
    emoji = _topic_emoji(title, content)

    # ── Header roxo com logo ──────────────────────────────────────────────────
    logo_rl = None
    logo_p  = Path(LOGO_PATH)
    if not logo_p.exists():
        for alt in ["assets/tati.png","assets/professor.jpg","assets/tati.jpg"]:
            if Path(alt).exists(): logo_p = Path(alt); break
    if logo_p.exists():
        try:
            circ   = _circle_crop(logo_p.read_bytes(), 72)
            logo_rl = RLImage(io.BytesIO(circ), width=1.3*cm, height=1.3*cm)
        except Exception:
            logo_rl = None

    t_para = Paragraph(f"{emoji}  {title}", s_htitle)
    s_para = Paragraph("✦ Tati's English Class ✦", s_hsub)

    if logo_rl:
        hd = Table([[logo_rl, [t_para, Spacer(1,4), s_para]]],
                   colWidths=[1.7*cm, W-1.7*cm])
    else:
        hd = Table([[[t_para, Spacer(1,4), s_para]]], colWidths=[W])

    hd.setStyle(TableStyle([
        ("BACKGROUND",   (0,0),(-1,-1), C_PURPLE),
        ("VALIGN",       (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",   (0,0),(-1,-1), 16), ("BOTTOMPADDING",(0,0),(-1,-1), 16),
        ("LEFTPADDING",  (0,0),(-1,-1), 12), ("RIGHTPADDING", (0,0),(-1,-1), 12),
        ("ROWBACKGROUNDS",(0,0),(-1,-1),[C_PURPLE]),
    ]))
    story.append(hd)
    story.append(ColorBar(C_GOLD, 4))
    story.append(Spacer(1, 10))

    # ── Palavras-chave ────────────────────────────────────────────────────────
    kws = _keywords(title, content)
    if kws:
        story.append(_box(
            Paragraph(f"<b>🔑 Key Vocabulary & Signal Words:</b><br/>{'  ·  '.join(kws)}", s_kw),
            C_PURPLE3, colors.HexColor("#c4b5fd"), pv=10
        ))
        story.append(Spacer(1, 8))

    # ── Conteúdo ──────────────────────────────────────────────────────────────
    tip_buf=[]; gram_buf=[]; ans_buf=[]; in_ans=False

    def flush_tip():
        if not tip_buf: return
        story.append(_box(Paragraph("💡 <b>Teacher's Tip:</b> "+" ".join(tip_buf), s_tip),
                          C_TIP_BG, colors.HexColor("#fbbf24"), pv=10))
        story.append(Spacer(1,6)); tip_buf.clear()

    def flush_gram():
        if not gram_buf: return
        story.append(_box(Paragraph("<br/>".join(gram_buf), s_gram),
                          C_BLUE_L, colors.HexColor("#93c5fd"), pv=10))
        story.append(Spacer(1,6)); gram_buf.clear()

    def flush_ans():
        if not ans_buf: return
        story.append(_box(Paragraph("<br/>".join(ans_buf), s_ansi),
                          C_ANS_BG, colors.HexColor("#6ee7b7"), pv=10))
        story.append(Spacer(1,4)); ans_buf.clear()

    for typ, tx in _parse_content(content):
        if typ == "blank":
            flush_tip(); flush_gram(); story.append(Spacer(1,4))
        elif typ == "divider":
            flush_tip(); flush_gram()
            story.append(HRFlowable(width="100%", thickness=0.5, color=C_BORDER, spaceAfter=4, spaceBefore=4))
        elif typ == "subtitle":
            flush_tip(); flush_gram()
            story.append(Paragraph(tx, s_sec))
            story.append(ColorBar(C_PURPLE, 2)); story.append(Spacer(1,6)); in_ans=False
        elif typ == "exheader":
            flush_tip(); flush_gram(); flush_ans()
            xt = Table([[Paragraph(tx or "Exercise", s_exh)]], colWidths=[W])
            xt.setStyle(TableStyle([
                ("BACKGROUND",(0,0),(-1,-1),C_PURPLE),
                ("TOPPADDING",(0,0),(-1,-1),8),("BOTTOMPADDING",(0,0),(-1,-1),8),
                ("LEFTPADDING",(0,0),(-1,-1),12),("RIGHTPADDING",(0,0),(-1,-1),12),
                ("ROWBACKGROUNDS",(0,0),(-1,-1),[C_PURPLE]),
            ]))
            story.append(Spacer(1,8)); story.append(xt); story.append(Spacer(1,6)); in_ans=False
        elif typ == "anskey":
            flush_tip(); flush_gram(); flush_ans()
            story.append(Spacer(1,8))
            story.append(HRFlowable(width="100%",thickness=1,color=C_GREEN,spaceAfter=4,spaceBefore=4))
            story.append(Paragraph(tx or "Answer Key", s_ans)); in_ans=True
        elif typ == "tip":
            flush_gram();
            if tx: tip_buf.append(tx)
        elif typ == "grammar":
            flush_tip();
            if tx: gram_buf.append(tx)
        elif typ == "exitem":
            flush_tip(); flush_gram()
            if in_ans: ans_buf.append(tx)
            else:      story.append(Paragraph(tx, s_exnum))
        elif typ == "body":
            flush_tip(); flush_gram()
            if in_ans: ans_buf.append(tx)
            else:      story.append(Paragraph(tx, s_body))

    flush_tip(); flush_gram(); flush_ans()

    # ── Dica pedagógica + rodapé ──────────────────────────────────────────────
    story.append(Spacer(1,14))
    story.append(HRFlowable(width="100%", thickness=0.5, color=C_BORDER))
    story.append(Spacer(1,6))
    story.append(_box(
        Paragraph(f"<b>📝 Teacher's Note:</b> {_pedagogical_tip(title,content)}",
                  S("TF", fontSize=9, fontName="Helvetica-Oblique",
                    textColor=colors.HexColor("#6b7280"), leading=14)),
        C_TIP_BG, colors.HexColor("#fbbf24"), pv=8
    ))
    story.append(Spacer(1,8))
    story.append(Paragraph("© Tati's English Class · All rights reserved", s_foot))
    doc.build(story)


# ══════════════════════════════════════════════════════════════════════════════
# GERADOR DOCX
# ══════════════════════════════════════════════════════════════════════════════

def _gen_docx(title: str, content: str, out_path: Path) -> None:
    from docx           import Document
    from docx.shared    import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns   import qn
    from docx.oxml      import OxmlElement

    def _set_cell_bg(cell, hex_color: str):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd  = OxmlElement("w:shd")
        shd.set(qn("w:val"),   "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"),  hex_color.lstrip("#"))
        tcPr.append(shd)

    doc = Document()
    for sec in doc.sections:
        sec.top_margin = sec.bottom_margin = sec.left_margin = sec.right_margin = Cm(2.0)

    emoji = _topic_emoji(title, content)

    # ── Header em tabela com fundo roxo ───────────────────────────────────────
    hdr_tbl = doc.add_table(rows=1, cols=1)
    hdr_tbl.style = "Table Grid"
    cell = hdr_tbl.cell(0, 0)
    _set_cell_bg(cell, "7c3aed")
    cell.paragraphs[0].clear()

    p1 = cell.paragraphs[0]
    p1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r1 = p1.add_run(f"{emoji}  {title}")
    r1.font.bold = True; r1.font.size = Pt(18)
    r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = cell.add_paragraph("✦ Tati's English Class ✦")
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.runs[0]; r2.font.size = Pt(9)
    r2.font.color.rgb = RGBColor(0xA7, 0x8B, 0xFA)

    # faixa dourada (linha horizontal)
    doc.add_paragraph()
    hr = doc.add_paragraph("─" * 80)
    hr.runs[0].font.color.rgb = RGBColor(0xF0, 0xA5, 0x00)
    hr.runs[0].font.size = Pt(8)

    # ── Palavras-chave ────────────────────────────────────────────────────────
    kws = _keywords(title, content)
    if kws:
        kw_tbl = doc.add_table(rows=1, cols=1)
        kw_tbl.style = "Table Grid"
        kc = kw_tbl.cell(0, 0)
        _set_cell_bg(kc, "f5f3ff")
        kp = kc.paragraphs[0]
        kr = kp.add_run("🔑 Key Vocabulary:  " + "  ·  ".join(kws))
        kr.font.size = Pt(9); kr.font.color.rgb = RGBColor(0x4C, 0x1D, 0x95)
        doc.add_paragraph()

    # ── Conteúdo ──────────────────────────────────────────────────────────────
    in_ans = False
    for typ, tx in _parse_content(content):
        if typ == "blank":
            doc.add_paragraph()
        elif typ == "divider":
            doc.add_paragraph("─" * 60)
        elif typ == "subtitle":
            p = doc.add_heading(tx, level=2)
            for r in p.runs: r.font.color.rgb = RGBColor(0x7C, 0x3A, 0xED)
            in_ans = False
        elif typ == "exheader":
            p = doc.add_heading(tx or "Exercise", level=3)
            for r in p.runs: r.font.color.rgb = RGBColor(0x7C, 0x3A, 0xED)
            in_ans = False
        elif typ == "anskey":
            p = doc.add_heading(tx or "Answer Key", level=2)
            for r in p.runs: r.font.color.rgb = RGBColor(0x05, 0x96, 0x69)
            in_ans = True
        elif typ in ("tip", "grammar"):
            p = doc.add_paragraph()
            r = p.add_run(("💡 Tip: " if typ == "tip" else "📘 ") + tx)
            r.font.italic = True; r.font.size = Pt(10)
            r.font.color.rgb = RGBColor(0x92, 0x40, 0x0E) if typ == "tip" else RGBColor(0x1E, 0x3A, 0x5F)
        elif typ == "exitem":
            p = doc.add_paragraph(tx, style="List Number" if not in_ans else "Normal")
            p.runs[0].font.size = Pt(11)
        elif typ == "body":
            p = doc.add_paragraph(tx)
            p.runs[0].font.size = Pt(11) if p.runs else None

    # ── Nota pedagógica ───────────────────────────────────────────────────────
    doc.add_paragraph()
    note_tbl = doc.add_table(rows=1, cols=1)
    note_tbl.style = "Table Grid"
    nc = note_tbl.cell(0, 0)
    _set_cell_bg(nc, "fffbeb")
    np_ = nc.paragraphs[0]
    nr  = np_.add_run("📝 Teacher's Note: " + _pedagogical_tip(title, content))
    nr.font.italic = True; nr.font.size = Pt(9)
    nr.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    doc.add_paragraph()
    fp = doc.add_paragraph("© Tati's English Class · All rights reserved")
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fp.runs[0].font.size = Pt(8); fp.runs[0].font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)

    doc.save(str(out_path))


# ══════════════════════════════════════════════════════════════════════════════
# GERADOR PPTX
# ══════════════════════════════════════════════════════════════════════════════

def _gen_pptx(title: str, content: str, out_path: Path) -> None:
    from pptx                  import Presentation
    from pptx.util             import Inches, Pt, Emu
    from pptx.dml.color        import RGBColor as PptxRGB
    from pptx.enum.text        import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    PURPLE = PptxRGB(0x7C, 0x3A, 0xED)
    GOLD   = PptxRGB(0xF0, 0xA5, 0x00)
    WHITE  = PptxRGB(0xFF, 0xFF, 0xFF)
    DARK   = PptxRGB(0x0D, 0x11, 0x17)
    PURPLE2= PptxRGB(0xC4, 0xB5, 0xFD)

    blank = prs.slide_layouts[6]  # blank

    def _add_bg(slide, color: PptxRGB):
        from pptx.oxml.ns import qn
        from lxml import etree
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _txt_box(slide, text, l, t, w, h, sz=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
        txb  = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf   = txb.text_frame; tf.word_wrap = wrap
        p    = tf.paragraphs[0]; p.alignment = align
        run  = p.add_run(); run.text = text
        run.font.size  = Pt(sz); run.font.bold = bold
        run.font.color.rgb = color
        return txb

    def _rect(slide, l, t, w, h, fill_color: PptxRGB):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    emoji = _topic_emoji(title, content)
    kws   = _keywords(title, content)

    # ── Slide 1: capa ─────────────────────────────────────────────────────────
    sl1 = prs.slides.add_slide(blank)
    _add_bg(sl1, PptxRGB(0x0D, 0x11, 0x17))
    _rect(sl1, 0, 0, 13.33, 7.5, PURPLE)
    _rect(sl1, 0, 6.9, 13.33, 0.6, GOLD)
    _txt_box(sl1, f"{emoji}  {title}", 0.5, 2.0, 12.33, 2.0,
             sz=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt_box(sl1, "✦ Tati's English Class ✦", 0.5, 4.2, 12.33, 0.8,
             sz=18, color=PURPLE2, align=PP_ALIGN.CENTER)
    if kws:
        kw_line = "  ·  ".join(kws[:6])
        _txt_box(sl1, f"🔑 {kw_line}", 0.5, 5.2, 12.33, 0.8,
                 sz=14, color=PptxRGB(0xC4,0xB5,0xFD), align=PP_ALIGN.CENTER)

    # ── Slides de conteúdo ────────────────────────────────────────────────────
    parsed   = _parse_content(content)
    cur_slide= None
    cur_y    = 0.8
    cur_texts= []

    def _flush_slide(texts, slide_title):
        if not texts or cur_slide is None: return
        y = 1.4
        for typ2, tx2 in texts:
            if y > 6.8: break
            if typ2 in ("body","exitem","grammar","tip"):
                prefix = "💡 " if typ2=="tip" else ("📘 " if typ2=="grammar" else "")
                c2 = PptxRGB(0x92,0x40,0x0E) if typ2=="tip" else (
                     PptxRGB(0x1E,0x3A,0x5F) if typ2=="grammar" else DARK)
                _txt_box(cur_slide, prefix+tx2, 0.5, y, 12.0, 0.7, sz=16, color=c2)
                y += 0.75
            elif typ2 == "anskey":
                _txt_box(cur_slide, "✅ "+tx2, 0.5, y, 12.0, 0.5, sz=14, bold=True,
                         color=PptxRGB(0x05,0x96,0x69))
                y += 0.6

    for typ, tx in parsed:
        if typ in ("subtitle","exheader"):
            # cria novo slide para cada seção
            cur_slide = prs.slides.add_slide(blank)
            _add_bg(cur_slide, PptxRGB(0xF9,0xFA,0xFB))
            _rect(cur_slide, 0, 0, 13.33, 1.1, PURPLE)
            _rect(cur_slide, 0, 1.1, 13.33, 0.06, GOLD)
            _txt_box(cur_slide, tx, 0.5, 0.15, 12.33, 0.8,
                     sz=24, bold=True, color=WHITE)
            cur_y = 1.4; cur_texts = []
        elif typ not in ("blank","divider") and cur_slide is not None:
            cur_texts.append((typ, tx))
            y = cur_y
            if typ in ("body","exitem","grammar","tip"):
                prefix = "💡 " if typ=="tip" else ("📘 " if typ=="grammar" else "")
                col = PptxRGB(0x92,0x40,0x0E) if typ=="tip" else (
                      PptxRGB(0x1E,0x3A,0x5F) if typ=="grammar" else DARK)
                _txt_box(cur_slide, prefix+tx, 0.5, cur_y, 12.0, 0.7, sz=16, color=col)
                cur_y += 0.78
            elif typ == "anskey":
                _txt_box(cur_slide, "✅ "+tx, 0.5, cur_y, 12.0, 0.5, sz=14, bold=True,
                         color=PptxRGB(0x05,0x96,0x69))
                cur_y += 0.62
            if cur_y > 6.6:
                cur_slide = prs.slides.add_slide(blank)
                _add_bg(cur_slide, PptxRGB(0xF9,0xFA,0xFB))
                _rect(cur_slide, 0, 0, 13.33, 0.5, PURPLE)
                _txt_box(cur_slide, "(cont.)", 0.5, 0.05, 12.33, 0.4, sz=14, color=WHITE)
                cur_y = 0.7

    # ── Slide final ───────────────────────────────────────────────────────────
    sl_end = prs.slides.add_slide(blank)
    _add_bg(sl_end, PptxRGB(0x0D,0x11,0x17))
    _rect(sl_end, 0, 0, 13.33, 7.5, PURPLE)
    tip_text = _pedagogical_tip(title, content)
    _txt_box(sl_end, "📝 Teacher's Note", 0.5, 1.5, 12.33, 0.8,
             sz=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _txt_box(sl_end, tip_text, 0.5, 2.6, 12.33, 2.5, sz=18, color=WHITE, align=PP_ALIGN.CENTER)
    _txt_box(sl_end, "© Tati's English Class", 0.5, 6.8, 12.33, 0.5,
             sz=12, color=PURPLE2, align=PP_ALIGN.CENTER)

    prs.save(str(out_path))


# ══════════════════════════════════════════════════════════════════════════════
# GERADOR TXT
# ══════════════════════════════════════════════════════════════════════════════

def _gen_txt(title: str, content: str, out_path: Path) -> None:
    lines = [
        f"{'='*60}",
        f"  {title.upper()}",
        f"  Tati's English Class",
        f"{'='*60}",
        "",
    ]
    kws = _keywords(title, content)
    if kws:
        lines += ["KEY VOCABULARY:", "  " + "  ·  ".join(kws), ""]

    for typ, tx in _parse_content(content):
        if typ == "blank":     lines.append("")
        elif typ == "divider": lines.append("-"*40)
        elif typ == "subtitle":lines += ["", f"[ {tx.upper()} ]", ""]
        elif typ == "exheader":lines += ["", f">> {tx} <<", ""]
        elif typ == "anskey":  lines += ["", "=== ANSWER KEY ===", ""]
        elif typ == "tip":     lines.append(f"💡 Tip: {tx}")
        elif typ == "grammar": lines.append(f"📘 {tx}")
        elif typ == "exitem":  lines.append(f"  {tx}")
        elif typ == "body":    lines.append(tx)

    lines += [
        "", "-"*40,
        f"Teacher's Note: {_pedagogical_tip(title, content)}",
        "", "© Tati's English Class",
    ]
    out_path.write_text("\n".join(lines), encoding="utf-8")


# ══════════════════════════════════════════════════════════════════════════════
# HELPERS COMUNS
# ══════════════════════════════════════════════════════════════════════════════

_TOPIC_EMOJIS = {
    "cafe":"☕","coffee":"☕","food":"🍽️","restaurant":"🍽️",
    "present perfect":"⏳","past simple":"📅","future":"🔮",
    "modal":"🔑","vocabulary":"📖","song":"🎵","music":"🎵",
    "travel":"✈️","business":"💼","sport":"⚽","conditional":"🌿",
    "pronunc":"🗣️","grammar":"📗","writing":"✍️","reading":"📖",
}

def _topic_emoji(title: str, content: str) -> str:
    m = (title + " " + content).lower()
    for k, v in _TOPIC_EMOJIS.items():
        if k in m: return v
    return "📚"

_KW_MAP = {
    "present perfect":["have/has + past participle","ever","never","already","yet","just","for","since","recently"],
    "past simple":    ["regular verbs (-ed)","irregular verbs","ago","last","yesterday","in + year","when"],
    "modal":          ["can","could","should","must","might","would","shall","may","have to"],
    "vocabulary":     ["noun","verb","adjective","adverb","synonym","antonym","context clue","collocations"],
    "food":           ["order","menu","bill","waiter","appetizer","main course","dessert","tip","portion"],
    "cafe":           ["coffee","espresso","latte","cappuccino","pastry","table","check","barista","to go"],
    "travel":         ["passport","boarding pass","customs","itinerary","accommodation","departure","arrival"],
    "business":       ["negotiate","agenda","deadline","stakeholder","proposal","budget","KPI","ROI"],
    "song":           ["lyrics","chorus","verse","bridge","rhythm","melody","metaphor","theme","imagery"],
    "conditional":    ["if clause","result clause","hypothetical","real","unreal","mixed","unless","provided"],
    "pronunc":        ["stress","intonation","vowel","consonant","minimal pair","linking","schwa"],
}

def _keywords(title: str, content: str) -> list:
    m = (title + " " + content).lower()
    for k, words in _KW_MAP.items():
        if k in m: return words[:8]
    return []

def _pedagogical_tip(title: str, content: str) -> str:
    t = (title + " " + content).lower()
    if "present perfect" in t:
        return ("Focus on the time connection — Present Perfect links PAST to NOW. "
                "Key signals: ever, never, already, yet, just, for, since.")
    if "past simple" in t:
        return "Contrast regular (-ed) and irregular verbs. Time markers: yesterday, last week, ago."
    if "modal" in t:
        return "Modals never take -s in 3rd person and always use the base form after them."
    if "song" in t or "music" in t:
        return "Listen 3 times: 1st for feeling, 2nd for keywords, 3rd sing along!"
    if "vocabulary" in t:
        return "Spaced repetition: review after 1 day, 1 week, 1 month. Use words in sentences."
    if "cafe" in t or "coffee" in t or "restaurant" in t:
        return "Role-play the dialogue! Real-life contexts make vocabulary stick 3× faster."
    if "conditional" in t:
        return "2nd conditional = hypothetical: If + Past Simple → would + base form."
    return ("Complete all exercises before checking the Answer Key. "
            "The retrieval effort is what builds lasting memory!")

def _parse_content(raw: str) -> list:
    out = []
    for line in raw.replace("\\n", "\n").split("\n"):
        s = line.strip()
        if not s:
            out.append(("blank",""))
        elif re.match(r"===\s*GRAMMAR", s, re.I):
            out.append(("subtitle","Grammar Explanation"))
        elif re.match(r"===\s*VOCAB", s, re.I):
            out.append(("subtitle","Vocabulary"))
        elif re.match(r"===\s*EXERCISE\s*\d+", s, re.I):
            n = re.search(r"\d+", s)
            out.append(("exheader", f"Exercise {n.group() if n else ''}"))
        elif re.match(r"===\s*ANSWER\s*KEY", s, re.I):
            out.append(("anskey","Answer Key"))
        elif re.match(r"===\s*TIP", s, re.I):
            out.append(("tip",""))
        elif re.match(r"^(TIP:|NOTE:|REMEMBER:)", s, re.I):
            out.append(("tip", s))
        elif re.match(r"^(GRAMMAR:|FORM:|STRUCTURE:|EXAMPLES?:)", s, re.I):
            out.append(("grammar", s))
        elif re.match(r"^(Exercise\s*\d+|EXERCISE\s*\d+)[\s:–-]", s):
            out.append(("exheader", s))
        elif re.match(r"^(Answer\s*Key|ANSWER\s*KEY)", s, re.I):
            out.append(("anskey", s))
        elif re.match(r"^\d+[\.\)]\s", s):
            out.append(("exitem", s))
        elif re.match(r"^[-•·]\s", s):
            out.append(("exitem", s[2:].strip()))
        elif re.match(r"^---+$", s):
            out.append(("divider",""))
        else:
            out.append(("body", s))
    return out


# ══════════════════════════════════════════════════════════════════════════════
# LEITURA DE ARQUIVOS ENVIADOS
# ══════════════════════════════════════════════════════════════════════════════

AUDIO_EXTS = {".mp3",".wav",".ogg",".m4a",".webm",".flac"}
IMAGE_EXTS = {".png",".jpg",".jpeg",".webp"}

def extract_file(raw: bytes, filename: str) -> dict:
    suffix = Path(filename).suffix.lower()
    if suffix in AUDIO_EXTS:
        return {"kind":"audio","label":"Áudio"}
    if suffix in IMAGE_EXTS:
        mm = {".png":"image/png",".jpg":"image/jpeg",".jpeg":"image/jpeg",".webp":"image/webp"}
        return {"kind":"image","label":"Imagem",
                "b64":base64.b64encode(raw).decode(),"media_type":mm.get(suffix,"image/jpeg")}
    if suffix == ".pdf":
        return {"kind":"text","label":"PDF","text":_read_pdf(raw)}
    if suffix in {".docx",".doc"}:
        return {"kind":"text","label":"Word","text":_read_docx(raw)}
    if suffix == ".txt":
        return {"kind":"text","label":"Texto","text":raw.decode("utf-8",errors="replace")}
    return {"kind":"unknown","label":suffix or "Arquivo"}

def _read_pdf(raw: bytes) -> str:
    try:
        import pdfplumber, io as _io
        parts = []
        with pdfplumber.open(_io.BytesIO(raw)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t: parts.append(t)
        return "\n\n".join(parts).strip() or "⚠️ Nenhum texto encontrado no PDF."
    except Exception as e:
        return f"❌ Erro ao extrair PDF: {e}"

def _read_docx(raw: bytes) -> str:
    try:
        import docx, io as _io
        d = docx.Document(_io.BytesIO(raw))
        return "\n\n".join(p.text for p in d.paragraphs if p.text.strip()) or "⚠️ Documento vazio."
    except Exception as e:
        return f"❌ Erro ao extrair DOCX: {e}"